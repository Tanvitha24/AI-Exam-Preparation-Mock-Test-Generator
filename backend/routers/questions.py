from openai import OpenAI
from dotenv import load_dotenv
from supabase import Client
from db import supabase
import io
import PyPDF2
from docx import Document
from pptx import Presentation
try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

from fastapi import APIRouter, Depends, HTTPException, status, Body, UploadFile, File
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os

# Load environment variables
load_dotenv()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

router = APIRouter()

class QuestionCreate(BaseModel):
    test_id: str
    content: str
    options: List[str]
    correct_answer: str
    explanation: Optional[str] = None
    topic: str
    difficulty: str
    question_type: str

class AIQuestionGenerate(BaseModel):
    topics: List[str]
    difficulty: str
    count: int
    question_types: List[str]

class DocumentQuestionGenerate(BaseModel):
    difficulty: str
    count: int
    question_types: List[str]

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from PDF: {str(e)}"
        )

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word document"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from Word document: {str(e)}"
        )

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from PowerPoint: {str(e)}"
        )

def extract_text_from_excel(file_content: bytes, file_extension: str) -> str:
    """Extract text from Excel file"""
    if pd is None:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Excel support requires pandas library. Please install it: pip install pandas openpyxl"
        )
    try:
        excel_file = io.BytesIO(file_content)
        if file_extension == 'xlsx':
            df = pd.read_excel(excel_file, engine='openpyxl')
        else:
            # Try xlrd, fallback to openpyxl if not available
            try:
                df = pd.read_excel(excel_file, engine='xlrd')
            except:
                df = pd.read_excel(excel_file, engine='openpyxl')
        
        text = ""
        for column in df.columns:
            text += f"{column}: "
            text += " ".join(df[column].astype(str).dropna().tolist()) + "\n"
        return text
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from Excel: {str(e)}"
        )

def extract_text_from_image(file_content: bytes) -> str:
    """Extract text from image using OCR"""
    if Image is None or pytesseract is None:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Image OCR support requires Pillow and pytesseract. Please install: pip install Pillow pytesseract. Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract"
        )
    try:
        image = Image.open(io.BytesIO(file_content))
        # Use pytesseract for OCR
        text = pytesseract.image_to_string(image)
        if not text or len(text.strip()) < 10:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Could not extract text from image. The image may not contain readable text or OCR failed."
            )
        return text
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from image using OCR: {str(e)}. Please ensure tesseract-ocr is installed on your system."
        )

def extract_text_from_document(file: UploadFile, file_content: bytes) -> str:
    """Extract text from document based on file type"""
    file_extension = file.filename.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_extension in ['doc', 'docx']:
        return extract_text_from_docx(file_content)
    elif file_extension in ['ppt', 'pptx']:
        return extract_text_from_pptx(file_content)
    elif file_extension in ['xls', 'xlsx']:
        return extract_text_from_excel(file_content, file_extension)
    elif file_extension in ['jpg', 'jpeg', 'png']:
        return extract_text_from_image(file_content)
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file type: {file_extension}. Supported types: PDF, DOC, DOCX, PPT, PPTX, XLS, XLSX, JPG, PNG"
        )

@router.post("/generate-from-document")
async def generate_questions_from_document(
    file: UploadFile = File(...),
    difficulty: str = Body(...),
    count: int = Body(...),
    question_types: str = Body(...)  # Receive as JSON string, parse it
):
    """Generate questions from uploaded document (PDF, PPT, Word)"""
    start_time = time.time()
    try:
        import json
        # Parse question_types from JSON string
        try:
            question_types_list = json.loads(question_types)
        except:
            question_types_list = ['multiple_choice', 'true_false']  # Default
        
        # Read file content
        file_content = await file.read()
        
        # Extract text from document
        extraction_start = time.time()
        document_text = extract_text_from_document(file, file_content)
        extraction_end = time.time()
        print(f"Text extraction took: {extraction_end - extraction_start:.2f} seconds")

        if not document_text or len(document_text.strip()) < 100:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Document appears to be empty or too short. Please upload a document with sufficient content."
            )
        
        # Truncate text if too long (OpenAI has token limits)
        max_chars = 10000  # Adjust based on your needs
        if len(document_text) > max_chars:
            document_text = document_text[:max_chars] + "... [truncated]"
        
        # Format the prompt based on question type
        question_type_instruction = ""
        if 'open_ended' in question_types_list:
            question_type_instruction = """
For open-ended questions, provide:
1. The question text
2. A sample answer or key points (as correct_answer)
3. An explanation of what a good answer should include
4. No options array needed
"""
        elif 'true_false' in question_types_list:
            question_type_instruction = """
For true/false questions, provide:
1. The question text
2. Options array with ["True", "False"]
3. The correct answer (either "True" or "False")
"""
        else:
            question_type_instruction = """
For multiple choice questions, provide:
1. The question text
2. 4 possible answers in options array
3. The correct answer
"""
        
        prompt = f"""Based on the following document content, generate {count} exam questions with {difficulty} difficulty.

Document Content:
{document_text}

{question_type_instruction}

For each question, provide:
1. The question text (as 'content')
2. Options array (for multiple choice/true-false) or omit for open-ended
3. The correct answer
4. A brief explanation
5. The question type (one of: {', '.join(question_types_list)})

The questions should be based on the content provided in the document above. Make sure the questions test understanding of the key concepts, facts, and information presented in the document.

Format as a JSON array of objects with fields: content, options (if applicable), correct_answer, explanation, topic, difficulty, question_type
"""
        
# Call Gemini API
        api_call_start = time.time()
        model = genai.GenerativeModel('gemini-pro')
        response = await model.generate_content_async(prompt)
        api_call_end = time.time()
        print(f"Gemini API call took: {api_call_end - api_call_start:.2f} seconds")
        
        # Extract the JSON string from the LLM response and parse it
        json_parsing_start = time.time()
        raw_content = response.text.strip()
        # Some models wrap JSON in markdown triple-backticks. Remove them if present
        if raw_content.startswith("```"):
            raw_content = raw_content.split("```", 2)[1]
            raw_content = raw_content.lstrip("json").strip()
        
        try:
            questions = json.loads(raw_content)
        except json.JSONDecodeError as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to parse generated questions JSON: {str(e)}"
            )
        json_parsing_end = time.time()
        print(f"JSON parsing took: {json_parsing_end - json_parsing_start:.2f} seconds")
        
        total_time = time.time() - start_time
        print(f"Total request time: {total_time:.2f} seconds")
        return {"questions": questions, "document_preview": document_text[:500]}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error in generate_questions_from_document: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Question generation from document failed: {str(e)}"
        )

@router.post("/generate")
async def generate_questions(request: AIQuestionGenerate):
    try:
        # Format the prompt for OpenAI
        prompt = f"""Generate {request.count} exam questions about {', '.join(request.topics)} with {request.difficulty} difficulty.
        For each question, provide:
        1. The question text
        2. 4 possible answers (for multiple choice)
        3. The correct answer
        4. A brief explanation
        5. The question type (one of: {', '.join(request.question_types)})
        
        Format as a JSON array of objects with fields: content, options, correct_answer, explanation, topic, difficulty, question_type
        """
        
# Call Gemini API
        model = genai.GenerativeModel('gemini-pro')
        response = await model.generate_content_async(prompt)
        
        # Extract the JSON string from the LLM response and parse it
        import json
        raw_content = response.text.strip()
        # Some models wrap JSON in markdown triple-backticks. Remove them if present
        if raw_content.startswith("```"):
            # remove first ``` and possible language tag
            raw_content = raw_content.split("```", 2)[1]
            # remove possible leading language hint like `json\n`
            raw_content = raw_content.lstrip("json").strip()
        
        try:
            questions = json.loads(raw_content)
        except json.JSONDecodeError as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to parse generated questions JSON: {str(e)}"
            )
        
        return {"questions": questions}
    except Exception as e:
        print(f"Error in generate_questions: {e}") # Added for debugging
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Question generation failed: {str(e)}"
        )

@router.post("/create")
async def create_question(question: QuestionCreate):
    try:
        response = supabase.table("questions").insert({
            "test_id": question.test_id,
            "content": question.content,
            "options": question.options,
            "correct_answer": question.correct_answer,
            "explanation": question.explanation,
            "topic": question.topic,
            "difficulty": question.difficulty,
            "question_type": question.question_type
        }).execute()
        
        if len(response.data) == 0:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Failed to create question"
            )
            
        return response.data[0]
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Question creation failed: {str(e)}"
        )

@router.get("/test/{test_id}")
async def get_test_questions(test_id: str):
    try:
        response = supabase.table("questions").select("*").eq("test_id", test_id).execute()
        return response.data
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to fetch questions: {str(e)}"
        )